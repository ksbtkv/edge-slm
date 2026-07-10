"""
PowerPoint ingestion for the Edge SLM pipeline.

Converts .pptx files into a section-based Document (one Section per slide).
Chunking is deferred to the source pack.
"""

from __future__ import annotations

import logging
import re
from datetime import date, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation

from ingestion.schema import (
    SOURCE_TYPE_LOCAL_FILE,
    Document,
    Section,
    new_document_id,
)


logger = logging.getLogger(__name__)


SUPPORTED_PPTX_EXTENSIONS = {".pptx"}


def clean_text(text: str) -> str:
    if not text:
        return ""

    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\t", " ")
    cleaned_lines = []

    for line in text.split("\n"):
        clean_line = re.sub(r"\s+", " ", line).strip()
        if clean_line:
            cleaned_lines.append(clean_line)

    return "\n".join(cleaned_lines).strip()


def safe_metadata_value(value: Any) -> Any:
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    return value


def extract_core_properties(presentation: Presentation) -> dict[str, Any]:
    core_properties = presentation.core_properties
    return {
        "title": safe_metadata_value(core_properties.title),
        "subject": safe_metadata_value(core_properties.subject),
        "author": safe_metadata_value(core_properties.author),
        "keywords": safe_metadata_value(core_properties.keywords),
        "comments": safe_metadata_value(core_properties.comments),
        "category": safe_metadata_value(core_properties.category),
        "created": safe_metadata_value(core_properties.created),
        "modified": safe_metadata_value(core_properties.modified),
        "last_modified_by": safe_metadata_value(core_properties.last_modified_by),
    }


def extract_text_from_text_frame(text_frame: Any) -> str:
    if text_frame is None:
        return ""
    return clean_text(text_frame.text)


def extract_text_from_table(shape: Any) -> str:
    if not getattr(shape, "has_table", False):
        return ""

    table_rows = []
    for row in shape.table.rows:
        cell_texts = [clean_text(cell.text) for cell in row.cells if clean_text(cell.text)]
        if cell_texts:
            table_rows.append(" | ".join(cell_texts))

    return "\n".join(table_rows).strip()


def extract_text_from_shape(shape: Any) -> list[str]:
    extracted_texts: list[str] = []

    if getattr(shape, "has_text_frame", False):
        text = extract_text_from_text_frame(shape.text_frame)
        if text:
            extracted_texts.append(text)

    if getattr(shape, "has_table", False):
        table_text = extract_text_from_table(shape)
        if table_text:
            extracted_texts.append(table_text)

    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            extracted_texts.extend(extract_text_from_shape(child_shape))

    return extracted_texts


def get_slide_title(slide: Any, fallback_title: str) -> str:
    title_shape = getattr(slide.shapes, "title", None)

    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title_text = clean_text(title_shape.text)
        if title_text:
            return title_text

    for shape in slide.shapes:
        for text in extract_text_from_shape(shape):
            first_line = text.split("\n")[0].strip()
            if first_line and len(first_line.split()) <= 15:
                return first_line

    return fallback_title


def extract_speaker_notes(slide: Any) -> str:
    try:
        notes_text = clean_text(slide.notes_slide.notes_text_frame.text)
    except Exception:
        return ""

    if not notes_text or notes_text.lower().strip() == "click to add notes":
        return ""

    return notes_text


def extract_slide_text(
    slide: Any,
    *,
    slide_number: int,
    include_notes: bool = True,
) -> dict[str, Any]:
    slide_title = get_slide_title(slide=slide, fallback_title=f"Slide {slide_number}")

    visible_text_parts = []
    for shape in slide.shapes:
        for text in extract_text_from_shape(shape):
            if text:
                visible_text_parts.append(text)

    visible_text = clean_text("\n".join(visible_text_parts))
    notes_text = extract_speaker_notes(slide) if include_notes else ""

    combined_parts = []
    if slide_title:
        combined_parts.append(f"Slide title: {slide_title}")
    if visible_text:
        combined_parts.append(visible_text)
    if notes_text:
        combined_parts.append(f"Speaker notes:\n{notes_text}")

    combined_text = clean_text("\n\n".join(combined_parts))

    return {
        "slide_number": slide_number,
        "slide_title": slide_title,
        "visible_text": visible_text,
        "notes_text": notes_text,
        "combined_text": combined_text,
    }


def detect_title_from_pptx(
    slides: list[dict[str, Any]],
    metadata_title: str | None,
    fallback_title: str,
) -> str:
    if metadata_title:
        clean_metadata_title = clean_text(metadata_title)
        if clean_metadata_title:
            return clean_metadata_title

    if slides:
        first_slide_title = clean_text(slides[0].get("slide_title", ""))
        if first_slide_title:
            return first_slide_title

    return fallback_title


def ingest_pptx(
    input_path: str | Path,
    *,
    title: str | None = None,
    language: str | None = "en",
    include_notes: bool = True,
) -> Document:
    """Ingest a .pptx file into the standardised Document schema."""
    pptx_path = Path(input_path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

    if not pptx_path.is_file():
        raise ValueError(f"Input path is not a file: {pptx_path}")

    if pptx_path.suffix.lower() not in SUPPORTED_PPTX_EXTENSIONS:
        raise ValueError(
            f"Unsupported PowerPoint file type: {pptx_path.suffix}. "
            f"Supported: {sorted(SUPPORTED_PPTX_EXTENSIONS)}"
        )

    presentation = Presentation(str(pptx_path))
    presentation_metadata = extract_core_properties(presentation)

    slides: list[dict[str, Any]] = []
    sections: list[Section] = []

    logger.info("Extracting PowerPoint slides: %s", pptx_path)

    for slide_index, slide in enumerate(presentation.slides):
        slide_number = slide_index + 1
        slide_info = extract_slide_text(
            slide=slide,
            slide_number=slide_number,
            include_notes=include_notes,
        )
        slides.append(slide_info)

        combined_text = slide_info["combined_text"]
        if not combined_text:
            continue

        sections.append(
            Section(
                index=len(sections),
                text=combined_text,
                raw_text=combined_text,
                heading=slide_info["slide_title"],
                slide_number=slide_number,
                extraction_method="python_pptx",
            )
        )

    if not sections:
        raise RuntimeError(f"No text content found in PowerPoint file: {pptx_path}")

    final_title = title or detect_title_from_pptx(
        slides=slides,
        metadata_title=presentation_metadata.get("title"),
        fallback_title=pptx_path.stem,
    )

    slides_with_notes = sum(1 for slide_info in slides if slide_info["notes_text"])

    return Document(
        document_id=new_document_id("pptx"),
        source_type=SOURCE_TYPE_LOCAL_FILE,
        source_path=str(pptx_path),
        modality="slides",
        content_type="slide_text",
        ingestor="pptx_ingestor",
        method="python_pptx",
        sections=sections,
        title=final_title,
        language=language,
        format_metadata={
            "pptx": {
                "title": presentation_metadata.get("title"),
                "author": presentation_metadata.get("author"),
                "subject": presentation_metadata.get("subject"),
                "keywords": presentation_metadata.get("keywords"),
                "comments": presentation_metadata.get("comments"),
                "category": presentation_metadata.get("category"),
                "created": presentation_metadata.get("created"),
                "modified": presentation_metadata.get("modified"),
                "last_modified_by": presentation_metadata.get("last_modified_by"),
                "slide_count": len(slides),
                "slides_with_text": len(sections),
                "slides_with_notes": slides_with_notes,
                "include_notes": include_notes,
            }
        },
    )
